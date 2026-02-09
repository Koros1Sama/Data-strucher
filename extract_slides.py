import os
import glob
from pptx import Presentation
import win32com.client
import sys

def clean_text(text):
    if not text:
        return ""
    # Replace Vertical Tab (VT) \x0b and Form Feed (FF) \x0c with newline
    text = text.replace('\x0b', '\n').replace('\x0c', '\n').replace('\r', '\n')
    return text

def extract_text_from_pptx(filepath):
    text_content = []
    try:
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            # Extract title if available
            if slide.shapes.title and slide.shapes.title.text:
                slide_text.append(f"--- Slide {i+1}: {clean_text(slide.shapes.title.text).strip()} ---")
            else:
                slide_text.append(f"--- Slide {i+1} ---")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Avoid duplicating title if it's already grabbed
                    if shape == slide.shapes.title:
                        continue
                    cleaned = clean_text(shape.text).strip()
                    if cleaned:
                        slide_text.append(cleaned)
            
            if slide_text:
                text_content.append("\n".join(slide_text))
                
    except Exception as e:
        return None, str(e)
    
    return "\n\n".join(text_content), None

def extract_text_from_ppt_com(filepath):
    # This function was not fully used in the main loop logic previously (it was inline), 
    # but let's keep it if we ever need to switch to function call. 
    # For now, I'll update the main loop where COM is used.
    pass 

def main():
    target_dir = r"D:\الجامعة\anti final\Data strucher\Data-strucher\Lessons\Old"
    
    # Files to process
    pptx_files = glob.glob(os.path.join(target_dir, "*.pptx"))
    ppt_files = glob.glob(os.path.join(target_dir, "*.ppt"))
    
    all_files = pptx_files + ppt_files
    
    if not all_files:
        print("No .ppt or .pptx files found in", target_dir)
        return

    print(f"Found {len(all_files)} files to process.")

    # Initialize COM only if needed and once
    ppt_app = None
    if ppt_files:
        try:
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        except Exception as e:
            print("Warning: Could not initialize PowerPoint. .ppt files might fail:", e)

    for filepath in all_files:
        filename = os.path.basename(filepath)
        print(f"Processing: {filename}...")
        
        ext = os.path.splitext(filename)[1].lower()
        output_path = filepath + ".txt"
        
        content = None
        error = None
        
        if ext == ".pptx":
            content, error = extract_text_from_pptx(filepath)
        elif ext == ".ppt":
            if ppt_app:
                try:
                    presentation = ppt_app.Presentations.Open(filepath, True, False, False)
                    content_list = []
                    for i, slide in enumerate(presentation.Slides):
                        # Try to get slide title safely
                        title = ""
                        try:
                            if slide.Shapes.Title.TextFrame.HasText:
                                title = clean_text(slide.Shapes.Title.TextFrame.TextRange.Text).strip()
                        except:
                            pass
                        
                        header = f"--- Slide {i+1}: {title} ---" if title else f"--- Slide {i+1} ---"
                        s_txt = [header]
                        
                        for shape in slide.Shapes:
                            try:
                                if shape.HasTextFrame and shape.TextFrame.HasText:
                                    raw_text = shape.TextFrame.TextRange.Text
                                    cleaned = clean_text(raw_text).strip()
                                    # Avoid title duplication
                                    if cleaned and cleaned != title:
                                        s_txt.append(cleaned)
                            except:
                                pass
                        content_list.append("\n".join(s_txt))
                    content = "\n\n".join(content_list)
                    presentation.Close()
                except Exception as e:
                    error = str(e)
            else:
                error = "PowerPoint application not available for .ppt file"

        if content:
            try:
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(content)
                print(f"Saved to: {os.path.basename(output_path)}")
            except Exception as e:
                print(f"Error writing text file for {filename}: {e}")
        else:
            print(f"Failed to extract text from {filename}. Error: {error}")

    if ppt_app:
        try:
            ppt_app.Quit()
        except:
            pass
    
    print("Done.")

if __name__ == "__main__":
    main()
